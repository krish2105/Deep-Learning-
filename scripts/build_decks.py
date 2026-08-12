#!/usr/bin/env python3
"""Generate both PowerPoint decks.

  deliverables/decks/SENTINEL-CXR_Week6_Progress.pptx   (10%, Week 6)
  deliverables/decks/SENTINEL-CXR_Final.pptx            (30%, Exam Week, 15 slides)

The visual language matches the application: a film-dark ground, one instrument
cyan used only for emphasis, and monospace for anything numeric. Run:
  python scripts/build_decks.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "deliverables" / "decks"
OUT.mkdir(parents=True, exist_ok=True)

# ── palette (mirrors apps/web/app/globals.css) ───────────────────────────
BASE = RGBColor(0x0B, 0x0D, 0x0E)
PANEL = RGBColor(0x16, 0x19, 0x1B)
SHOULDER = RGBColor(0x2A, 0x2F, 0x33)
MID = RGBColor(0x8A, 0x92, 0x99)
HIGHLIGHT = RGBColor(0xE8, 0xEC, 0xEF)
INSTRUMENT = RGBColor(0x2E, 0x9C, 0xB8)
STAT = RGBColor(0xD6, 0x45, 0x41)
URGENT = RGBColor(0xD9, 0x90, 0x3F)

W, H = Inches(13.333), Inches(7.5)
SANS, MONO = "Calibri", "Consolas"

TEAM = "Krishna Mathur AS25DXB018 · Atharva Soundankar AS25DXB020 · Yash Petkar AS25DXB021"


def deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BASE
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, x, y, w, h, content, *, size=18, bold=False, colour=HIGHLIGHT,
         font=SANS, align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP,
         italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = content.split("\n") if isinstance(content, str) else content
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        r = para.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = colour
    return box


def rule(slide, x, y, w, colour=SHOULDER, thickness=1.25):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Pt(thickness))
    ln.fill.solid()
    ln.fill.fore_color.rgb = colour
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def panel(slide, x, y, w, h, fill=PANEL, border=SHOULDER):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.03
    return sh


def eyebrow(slide, label, n=None, total=None):
    text(slide, 0.7, 0.45, 9, 0.35, label.upper(), size=10.5, colour=INSTRUMENT,
         font=MONO)
    if n is not None:
        text(slide, 11.4, 0.45, 1.3, 0.35, f"{n:02d}/{total}", size=10.5,
             colour=MID, font=MONO, align=PP_ALIGN.RIGHT)


def title(slide, t, y=1.0, size=40, colour=HIGHLIGHT):
    text(slide, 0.7, y, 11.9, 1.5, t, size=size, bold=True, colour=colour,
         spacing=0.92)


def footer(slide, note="SENTINEL-CXR · not a medical device"):
    text(slide, 0.7, 6.95, 11.9, 0.3, note, size=8.5, colour=MID, font=MONO)


def bullet_panels(slide, items, y=2.5, h=1.5, gap=0.22, cols=3, x0=0.7,
                  total_w=11.9):
    w = (total_w - gap * (cols - 1)) / cols
    for i, (head, body) in enumerate(items):
        col, row = i % cols, i // cols
        x = x0 + col * (w + gap)
        yy = y + row * (h + gap)
        panel(slide, x, yy, w, h)
        text(slide, x + 0.28, yy + 0.22, w - 0.56, 0.4, head, size=14, bold=True,
             colour=INSTRUMENT)
        text(slide, x + 0.28, yy + 0.68, w - 0.56, h - 0.9, body, size=11,
             colour=MID, spacing=1.05)


def table(slide, headers, rows, x=0.7, y=2.4, w=11.9, row_h=0.42,
          widths=None, size=11, highlight_row=None):
    n = len(headers)
    widths = widths or [w / n] * n
    # header
    cx = x
    for i, hcell in enumerate(headers):
        text(slide, cx, y, widths[i], row_h, hcell, size=size - 1, bold=True,
             colour=INSTRUMENT, font=MONO)
        cx += widths[i]
    rule(slide, x, y + row_h - 0.02, w, INSTRUMENT, 1.0)
    # rows
    for r, row in enumerate(rows):
        yy = y + row_h + 0.06 + r * row_h
        colour = HIGHLIGHT if highlight_row == r else MID
        bold = highlight_row == r
        cx = x
        for i, cell in enumerate(row):
            fnt = MONO if i > 0 else SANS
            text(slide, cx, yy, widths[i], row_h, str(cell), size=size,
                 colour=colour, font=fnt, bold=bold)
            cx += widths[i]
        if r < len(rows) - 1:
            rule(slide, x, yy + row_h - 0.04, w, SHOULDER, 0.5)


def chip(slide, x, y, w, label, value, chroma=1.0, hatched=False):
    """A confidence chip — saturation encodes certainty, as in the product."""
    panel(slide, x, y, w, 0.62)
    sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.18),
                                Inches(y + 0.21), Inches(0.2), Inches(0.2))
    sw.fill.solid()
    if hatched:
        sw.fill.fore_color.rgb = SHOULDER
    else:
        # Blend instrument -> mid by chroma, mirroring confidence_to_chroma.
        sw.fill.fore_color.rgb = RGBColor(
            int(MID[0] + (INSTRUMENT[0] - MID[0]) * chroma),
            int(MID[1] + (INSTRUMENT[1] - MID[1]) * chroma),
            int(MID[2] + (INSTRUMENT[2] - MID[2]) * chroma),
        )
    sw.line.fill.background()
    sw.shadow.inherit = False
    text(slide, x + 0.5, y + 0.15, w - 1.5, 0.35, label, size=12,
         colour=HIGHLIGHT if not hatched else MID)
    text(slide, x + w - 1.15, y + 0.15, 0.95, 0.35, value, size=12, font=MONO,
         colour=MID, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════
# TITLE / CLOSING (shared)
# ═════════════════════════════════════════════════════════════════════════
def title_slide(prs, kind: str, subtitle: str):
    s = blank(prs)
    rule(s, 0.7, 2.05, 3.0, INSTRUMENT, 2.0)
    text(s, 0.7, 1.45, 9, 0.4, "MAIB AI 114 · DEEP LEARNING", size=11,
         colour=INSTRUMENT, font=MONO)
    text(s, 0.7, 2.35, 11.9, 1.4, "SENTINEL-CXR", size=62, bold=True,
         colour=HIGHLIGHT)
    text(s, 0.7, 3.55, 11, 0.55, "Uncertainty-Aware Chest Radiograph Triage",
         size=22, colour=HIGHLIGHT)
    text(s, 0.7, 4.15, 11, 0.5, subtitle, size=14, colour=MID, italic=True)
    rule(s, 0.7, 5.15, 11.9)
    text(s, 0.7, 5.4, 11.9, 0.4, TEAM, size=12, colour=MID)
    text(s, 0.7, 5.85, 11.9, 0.4,
         "Prof Anshul Gupta · S P Jain School of Global Management, Dubai",
         size=11, colour=MID)
    text(s, 0.7, 6.3, 11.9, 0.4, kind, size=11, colour=INSTRUMENT, font=MONO)
    footer(s)
    return s


# ═════════════════════════════════════════════════════════════════════════
# DECK 1 — WEEK 6 PROGRESS (10%)
# ═════════════════════════════════════════════════════════════════════════
def build_week6():
    prs = deck()
    T = 10

    title_slide(prs, "GROUP PROJECT & CLASSROOM PRESENTATION · WEEK 6",
                "Progress review — problem, approach and current state")

    # 2 — the problem
    s = blank(prs); eyebrow(s, "the problem", 2, T)
    title(s, "Accuracy is not what blocks deployment.")
    text(s, 0.7, 2.5, 7.4, 2.6,
         "A model trained with cross-entropy emits a probability for every "
         "input — including a photograph, a mispositioned film, or a case where "
         "the evidence genuinely does not support a decision.\n\n"
         "Clinicians are measurably less likely to override a confident machine "
         "judgement. A system that is wrong and confident therefore causes more "
         "harm than one that says nothing.",
         size=15, colour=MID, spacing=1.2)
    panel(s, 8.5, 2.5, 4.1, 2.6)
    text(s, 8.8, 2.75, 3.5, 0.4, "WHAT WE BUILD INSTEAD", size=10,
         colour=INSTRUMENT, font=MONO)
    text(s, 8.8, 3.25, 3.5, 1.7,
         "A prediction set with a\ncoverage guarantee —\nand the discipline to\n"
         "abstain when that\nguarantee cannot be met.",
         size=14, colour=HIGHLIGHT, spacing=1.15)
    footer(s)

    # 3 — three defences
    s = blank(prs); eyebrow(s, "approach", 3, T)
    title(s, "Three layers of refusal.")
    bullet_panels(s, [
        ("01  Distributional gate",
         "A VAE trained only on radiographs. Reconstruction error rejects "
         "anything that is not a chest X-ray, before classification runs."),
        ("02  Epistemic uncertainty",
         "Separates the model's own ignorance from irreducible ambiguity in the "
         "image. Only ignorance justifies abstention."),
        ("03  Conformal prediction",
         "Converts uncalibrated scores into sets with distribution-free "
         "finite-sample coverage. Empty or oversized set means abstain."),
    ], y=2.5, h=2.4)
    footer(s)

    # 4 — dataset
    s = blank(prs); eyebrow(s, "data", 4, T)
    title(s, "One dataset, two structures.")
    text(s, 0.7, 2.15, 11.9, 0.5,
         "NIH ChestX-ray14 — 112,120 radiographs, 30,805 patients, 14 pathologies.",
         size=14, colour=MID)
    table(s, ["Column", "What it unlocks"], [
        ["Finding Labels", "14-label multi-label classification target"],
        ["Patient ID", "The splitting key — and patient timelines"],
        ["Follow-up #", "Sequence position: makes RNN/LSTM clinically real"],
        ["Age / Gender / View", "Fairness audit strata"],
    ], y=2.9, widths=[3.6, 8.3], row_h=0.55, size=13)
    panel(s, 0.7, 5.5, 11.9, 1.1)
    text(s, 1.0, 5.72, 11.3, 0.7,
         "We split by PATIENT, never by image. Splitting by image puts a patient's "
         "follow-up scans on both sides of the boundary — the model memorises the "
         "patient and every metric inflates.",
         size=13, colour=HIGHLIGHT, spacing=1.1)
    footer(s)

    # 5 — syllabus coverage
    s = blank(prs); eyebrow(s, "coverage", 5, T)
    title(s, "Every week of the unit, one system.")
    table(s, ["Wk", "Topic", "Where it lives in SENTINEL-CXR"], [
        ["1–2", "NN & backprop", "From-scratch NumPy backprop + gradient check"],
        ["3", "CNNs", "DenseNet-121 multi-label classifier"],
        ["4–5", "RNN / LSTM", "Progression over patient follow-up sequences"],
        ["6", "GANs", "DCGAN augmentation for rare pathologies"],
        ["7", "Autoencoders / VAE", "The out-of-distribution gate"],
        ["8", "Transfer learning", "Frozen vs full vs progressive unfreezing"],
        ["9", "Deep RL", "DQN worklist triage"],
        ["10", "ViT / CLIP", "CNN vs Transformer head-to-head"],
        ["11", "GenAI integration", "Grounded report generation"],
        ["12", "Ethics & fairness", "Disaggregated audit + model card"],
    ], y=2.25, widths=[0.9, 3.0, 8.0], row_h=0.4, size=12)
    footer(s)

    # 6 — architecture
    s = blank(prs); eyebrow(s, "architecture", 6, T)
    title(s, "Designed around a hard constraint.")
    text(s, 0.7, 2.15, 11.9, 0.5,
         "Render's free tier gives 512 MB. PyTorch does not fit. That asymmetry "
         "produced the whole design.", size=14, colour=MID)
    panel(s, 0.7, 2.9, 11.9, 2.9)
    text(s, 1.1, 3.15, 11.1, 2.5,
         "Vercel      Next.js 15 — landing + clinical console\n"
         "   |        HTTPS + JWT\n"
         "Render      FastAPI · 512 MB · NO PyTorch\n"
         "   |        conformal head (NumPy) · auth · audit log\n"
         "   +- cold -> ONNX int8 · <900 ms · 'reduced' mode\n"
         "   |\n"
         "HF Spaces   16 GB · DenseNet · ViT · VAE · LSTM · Grad-CAM\n"
         "   |\n"
         "Supabase    Postgres + object storage",
         size=13, colour=HIGHLIGHT, font=MONO, spacing=1.15)
    footer(s)

    # 7 — progress
    s = blank(prs); eyebrow(s, "progress", 7, T)
    title(s, "Where we are.")
    bullet_panels(s, [
        ("Done", "Full-stack system live on free-tier infrastructure. 83 tests "
                 "passing. Real inference in ~150 ms via a 7.9 MB quantised model "
                 "carried by the API itself."),
        ("Done", "Calibrated on 4,999 real radiographs: coverage 0.8845, and a "
                 "fairness audit that FAILS its own tolerance at 0.2149."),
        ("Next", "Group-trained weights to replace the pretrained checkpoint; "
                 "DQN training; VAE gate; wider calibration corpus."),
    ], y=2.5, h=2.3)
    footer(s)

    # 8 — what went wrong
    s = blank(prs); eyebrow(s, "findings so far", 8, T)
    title(s, "Three bugs worth reporting.")
    table(s, ["Found", "Why it mattered"], [
        ["Triage env could not separate policies",
         "Random, FIFO and heuristic scored identically — the benchmark was meaningless"],
        ["Pretrained checkpoint has no dropout",
         "MC-dropout would have reported epistemic uncertainty of exactly zero"],
        ["int8 model would not load in production",
         "Verified numerically but never verified as executable on the pinned runtime"],
    ], y=2.5, widths=[4.6, 7.3], row_h=0.85, size=13)
    panel(s, 0.7, 5.6, 11.9, 1.0)
    text(s, 1.0, 5.82, 11.3, 0.6,
         "Each was found by testing a claim the system makes about itself — not by "
         "checking whether the app looked right.",
         size=13, colour=HIGHLIGHT)
    footer(s)

    # 9 — plan
    s = blank(prs); eyebrow(s, "plan", 9, T)
    title(s, "To the final submission.")
    table(s, ["Stage", "Owner", "Output"], [
        ["Model, ONNX export, calibration", "Krishna Mathur", "conformal_calibration.json"],
        ["Deployment, resilience, CI gates", "Atharva Soundankar", "Live URLs, 83 tests"],
        ["Console, dashboard, fairness", "Yash Petkar", "fairness_report.json"],
        ["Report and both decks", "All three", "15 pp report, 25 slides"],
    ], y=2.6, widths=[5.2, 2.4, 4.3], row_h=0.62, size=13)
    footer(s)

    # 10 — close
    s = blank(prs); eyebrow(s, "thank you", 10, T)
    title(s, "Questions.", y=2.6, size=54)
    rule(s, 0.7, 3.9, 3.0, INSTRUMENT, 2.0)
    text(s, 0.7, 4.2, 11.9, 0.9,
         "github.com/krish2105/Deep-Learning-", size=15, colour=MID, font=MONO)
    text(s, 0.7, 5.0, 11.9, 0.5, TEAM, size=12, colour=MID)
    footer(s)

    path = OUT / "SENTINEL-CXR_Week6_Progress.pptx"
    prs.save(path)
    return path, len(prs.slides._sldIdLst)


# ═════════════════════════════════════════════════════════════════════════
# DECK 2 — FINAL (30%), 15 slides
# ═════════════════════════════════════════════════════════════════════════
def build_final():
    prs = deck()
    T = 15

    # 1
    title_slide(prs, "FINAL GROUP PROJECT · EXAM WEEK · 30%",
                "A diagnostic system that knows when it does not know")

    # 2 — thesis
    s = blank(prs); eyebrow(s, "thesis", 2, T)
    title(s, "It tells you when it doesn't know.", size=44)
    text(s, 0.7, 2.75, 8.0, 2.4,
         "Most medical deep learning fails to deploy not because it is "
         "insufficiently accurate, but because it is confident about everything "
         "— including inputs it has never seen.\n\n"
         "SENTINEL-CXR produces a prediction set with a statistical coverage "
         "guarantee, and abstains when that guarantee cannot be met.",
         size=16, colour=MID, spacing=1.2)
    chip(s, 8.9, 2.8, 3.7, "Effusion", "0.87", chroma=1.0)
    chip(s, 8.9, 3.55, 3.7, "Nodule", "0.61", chroma=0.35)
    chip(s, 8.9, 4.3, 3.7, "Abstained", "—", hatched=True)
    text(s, 8.9, 5.05, 3.7, 0.6, "Confidence is chroma.\nColour drains as the "
         "model doubts.", size=11, colour=MID, italic=True, spacing=1.1)
    footer(s)

    # 3 — problem
    s = blank(prs); eyebrow(s, "the problem", 3, T)
    title(s, "A confident wrong answer is worse than none.")
    bullet_panels(s, [
        ("Scores aren't probabilities",
         "Thresholding a cross-entropy output at 0.5 gives no guarantee about "
         "how often the answer is right."),
        ("Automation bias",
         "Clinicians are measurably less likely to override a confident machine "
         "judgement — so confidence must be earned."),
        ("Hidden stratification",
         "Aggregate AUROC stays high while the model fails badly on clinically "
         "important subgroups (Oakden-Rayner, 2020)."),
    ], y=2.5, h=2.4)
    footer(s)

    # 4 — architecture of refusal
    s = blank(prs); eyebrow(s, "method", 4, T)
    title(s, "Three layers of refusal.")
    bullet_panels(s, [
        ("01  VAE gate",
         "Trained only on radiographs. Reconstruction error rejects "
         "non-radiographs before classification. Threshold set at 1% FPR — an "
         "explicit design decision, not a round number."),
        ("02  Uncertainty",
         "Predictive entropy splits exactly into aleatoric (irreducible "
         "ambiguity) and epistemic (the model's ignorance). Only epistemic "
         "justifies abstention."),
        ("03  Conformal",
         "Nonconformity quantile at ceil((n+1)(1-α))/n. The finite-sample "
         "correction is what makes coverage exact rather than asymptotic."),
    ], y=2.5, h=2.6)
    footer(s)

    # 5 — data
    s = blank(prs); eyebrow(s, "data", 5, T)
    title(s, "Patient ID and Follow-up # change everything.")
    text(s, 0.7, 2.2, 11.9, 0.5,
         "NIH ChestX-ray14 · 112,120 radiographs · 30,805 patients · ~3–4 studies each",
         size=14, colour=MID, font=MONO)
    panel(s, 0.7, 2.95, 5.8, 2.7)
    text(s, 1.0, 3.2, 5.2, 0.4, "WHAT THEY UNLOCK", size=10, colour=INSTRUMENT, font=MONO)
    text(s, 1.0, 3.7, 5.2, 1.8,
         "Multiple studies per patient form real\nclinical timelines. The recurrent\n"
         "branch models disease progression\nrather than serving as an\nunmotivated exercise.",
         size=13, colour=HIGHLIGHT, spacing=1.15)
    panel(s, 6.8, 2.95, 5.8, 2.7)
    text(s, 7.1, 3.2, 5.2, 0.4, "AND WHAT THEY DEMAND", size=10, colour=STAT, font=MONO)
    text(s, 7.1, 3.7, 5.2, 1.8,
         "Split by PATIENT, never by image.\nSplitting by image leaks follow-up\n"
         "scans across the boundary and\ninflates every metric. This is the\n"
         "most common error on this dataset.",
         size=13, colour=HIGHLIGHT, spacing=1.15)
    footer(s)

    # 6 — architecture
    s = blank(prs); eyebrow(s, "architecture", 6, T)
    title(s, "512 MB shaped the whole system.")
    panel(s, 0.7, 2.3, 7.3, 3.6)
    text(s, 1.0, 2.55, 6.8, 3.2,
         "Vercel      Next.js 15\n"
         "   |        HTTPS + JWT\n"
         "Render      FastAPI · 512 MB · NO PyTorch\n"
         "   |        conformal head · audit log\n"
         "   +- cold -> ONNX int8 · <900 ms\n"
         "   |\n"
         "HF Spaces   16 GB · DenseNet · VAE\n"
         "   |        Grad-CAM · LSTM · ViT\n"
         "Supabase    Postgres + storage",
         size=13, colour=HIGHLIGHT, font=MONO, spacing=1.2)
    panel(s, 8.3, 2.3, 4.3, 3.6)
    text(s, 8.6, 2.55, 3.7, 0.4, "GRACEFUL DEGRADATION", size=10,
         colour=INSTRUMENT, font=MONO)
    text(s, 8.6, 3.05, 3.7, 2.7,
         "A free HF Space sleeps after\n48h and wakes in ~40s.\n\n"
         "When it is cold, the ONNX\nfast path answers immediately\n"
         "and the response is marked\nREDUCED — visibly, in the UI.\n\n"
         "A clinician must know which\nmode produced their result.",
         size=12, colour=MID, spacing=1.15)
    footer(s)

    # 7 — pipeline
    s = blank(prs); eyebrow(s, "pipeline", 7, T)
    title(s, "Eleven steps, one decision each.")
    panel(s, 0.7, 2.35, 11.9, 3.5)
    text(s, 1.1, 2.65, 11.1, 3.0,
         "1   upload            -> presigned URL, storage\n"
         "2   VAE OOD gate      -> not a radiograph?  REJECT, stop\n"
         "3   classifier        -> DenseNet-121, 14 sigmoid scores\n"
         "4   uncertainty       -> MC-dropout / TTA, epistemic vs aleatoric\n"
         "5   conformal head    -> prediction set at 90% coverage\n"
         "6   abstention        -> empty or oversized?  ROUTE TO HUMAN\n"
         "7   Grad-CAM          -> per-finding activation map\n"
         "8   progression       -> LSTM over prior visits, if any\n"
         "9   triage            -> DQN priority, worklist position\n"
         "10  report            -> LLM grounded strictly in steps 3–9\n"
         "11  persist           -> audit entry, stream to client",
         size=13, colour=HIGHLIGHT, font=MONO, spacing=1.16)
    footer(s)

    # 8 — grounding
    s = blank(prs); eyebrow(s, "generative ai · outcome d", 8, T)
    title(s, "The model never sees the image.")
    text(s, 0.7, 2.2, 11.9, 0.5,
         "A language model asked to describe a radiograph will invent findings — "
         "fluently, in clinical register.", size=14, colour=MID)
    bullet_panels(s, [
        ("Closed input",
         "It receives structured model output only. No pixels, so it cannot "
         "invent a finding from the image."),
        ("Closed vocabulary",
         "The prompt enumerates exactly which pathologies may be named."),
        ("Verified output",
         "Generated text is scanned for unsupported findings. One hit and the "
         "generation is discarded and logged."),
    ], y=2.9, h=1.9)
    panel(s, 0.7, 5.05, 11.9, 1.2)
    text(s, 1.0, 5.25, 11.3, 0.9,
         "Layer 3 is what makes 1 and 2 trustworthy. Prompt instructions are a "
         "request; verification is a guarantee. Even 'no pneumothorax is seen' is "
         "rejected — it implies the system looked, and if it didn't, that's false.",
         size=13, colour=HIGHLIGHT, spacing=1.1)
    footer(s)

    # 9 — results 1
    s = blank(prs); eyebrow(s, "results", 9, T)
    title(s, "The claim, and the number that tests it.")
    text(s, 0.7, 2.15, 11.9, 0.5,
         "Coverage is the project's central claim — so it is asserted in CI.",
         size=14, colour=MID)
    panel(s, 0.7, 2.9, 5.8, 2.5)
    text(s, 1.0, 3.15, 5.2, 0.4, "MEASURED COVERAGE", size=10,
         colour=INSTRUMENT, font=MONO)
    text(s, 1.0, 3.6, 5.2, 0.8, "0.8845", size=44, bold=True, colour=URGENT,
         font=MONO)
    text(s, 1.0, 4.5, 5.2, 0.9,
         "against a 0.90 target — BELOW\n4,999 radiographs, 1,335 patients\n"
         "patient-disjoint split",
         size=12, colour=MID, spacing=1.1)
    panel(s, 6.8, 2.9, 5.8, 2.5)
    text(s, 7.1, 3.15, 5.2, 0.4, "WHY IT FALLS SHORT", size=10,
         colour=INSTRUMENT, font=MONO)
    text(s, 7.1, 3.6, 5.2, 1.9,
         "Coverage assumes exchangeability.\nSplitting by PATIENT is required\n"
         "to avoid follow-up leakage — and\nbreaks exchangeability.\n\n"
         "Pneumonia has 31 calibration\npositives; its quantile is noisy.",
         size=12, colour=HIGHLIGHT, spacing=1.12)
    footer(s)

    # 10 — results 2
    s = blank(prs); eyebrow(s, "results", 10, T)
    title(s, "Triage: the RL environment had to be rebuilt.")
    table(s, ["Policy", "Mean return", "Std"], [
        ["Random", "-837.74", "41.45"],
        ["First-in-first-out", "-841.66", "31.48"],
        ["Urgency heuristic", "-397.73", "29.82"],
        ["Oracle-greedy (bound)", "-367.22", "26.08"],
    ], y=2.5, widths=[5.0, 3.4, 3.5], row_h=0.52, size=14, highlight_row=2)
    panel(s, 0.7, 4.95, 11.9, 1.4)
    text(s, 1.0, 5.15, 11.3, 1.1,
         "Our first environment paid a bonus for reading a study. It swamped the "
         "waiting cost, and random, FIFO and heuristic all scored within noise — "
         "the benchmark could not tell a good policy from a coin flip. Reformulating "
         "the reward as pure cost separated them and made the result meaningful.",
         size=13, colour=HIGHLIGHT, spacing=1.12)
    footer(s)

    # 11 — recurrent
    s = blank(prs); eyebrow(s, "results", 11, T)
    title(s, "Gated cells hold gradient. Which one wins is not settled.")
    table(s, ["Cell", "grad @ t=0", "grad @ t=59", "Parameters"], [
        ["Vanilla RNN", "1.46e-17", "1.09e-01", "463,374"],
        ["GRU", "8.69e-14", "5.61e-02", "1,382,926"],
        ["LSTM", "2.92e-14", "3.02e-02", "1,842,702"],
    ], y=2.5, widths=[3.2, 3.0, 3.0, 2.7], row_h=0.52, size=14)
    panel(s, 0.7, 4.7, 11.9, 1.6)
    text(s, 1.0, 4.92, 11.3, 1.3,
         "Gated cells retain ~4 orders of magnitude more gradient at t=0 over a "
         "60-step sequence, averaged across 12 seeds. GRU and LSTM are close and "
         "their ordering is NOT stable across initialisations — so we make no claim "
         "that either dominates on this probe. Parameter count is reported because "
         "it is a confound in any such comparison.",
         size=13, colour=MID, spacing=1.12)
    footer(s)

    # 12 — ethics
    s = blank(prs); eyebrow(s, "ethics · outcome e", 12, T)
    title(s, "The audit fails. We predicted why before we measured.")
    table(s, ["Stratum", "TPR gap", "FPR gap", "Within 0.10?"], [
        ["Patient sex", "0.0165", "0.0267", "yes"],
        ["Age band", "0.1295", "0.2040", "NO"],
        ["View position", "0.0228", "0.2149", "NO"],
    ], y=2.45, widths=[4.2, 2.6, 2.6, 2.5], row_h=0.5, size=14, highlight_row=2)
    panel(s, 0.7, 4.5, 11.9, 1.75)
    text(s, 1.0, 4.72, 11.3, 1.45,
         "Portable AP films are taken of patients too unwell to stand, so view "
         "position correlates with severity. The FPR gap is 0.2149 while the TPR "
         "gap is only 0.0228 — the model is not MISSING more disease on AP films, "
         "it is OVER-CALLING it. That asymmetry is the signature of a shortcut, "
         "and aggregate AUROC would never have shown it.",
         size=13, colour=HIGHLIGHT, spacing=1.14)
    text(s, 0.7, 6.4, 11.9, 0.4,
         "ChestX-ray14 has no race labels — that axis cannot be audited at all, "
         "and that absence is a finding, not an absence of bias.",
         size=12, bold=True, colour=STAT)
    footer(s)

    # 13 — limitations
    s = blank(prs); eyebrow(s, "limitations", 13, T)
    title(s, "What this system is not.")
    table(s, ["Limitation", "Consequence"], [
        ["Labels NLP-mined, ~10% error", "Every metric inherits that ceiling"],
        ["Single US institution", "Generalisation unvalidated"],
        ["Marginal, not simultaneous, coverage", "Per-label guarantee only"],
        ["Exchangeability assumed", "Distribution shift voids the guarantee"],
        ["Grad-CAM is correlational", "Not evidence of correct reasoning"],
        ["Triage results are simulated", "No claim about a real reading room"],
    ], y=2.45, widths=[6.2, 5.7], row_h=0.55, size=13)
    panel(s, 0.7, 6.0, 11.9, 0.75)
    text(s, 1.0, 6.16, 11.3, 0.5,
         "Not a medical device. No regulatory clearance. Must not inform patient care.",
         size=13, bold=True, colour=STAT)
    footer(s)

    # 14 — what we learned
    s = blank(prs); eyebrow(s, "reflection", 14, T)
    title(s, "The bugs taught us more than the metrics.")
    table(s, ["What broke", "How it was caught"], [
        ["Triage env couldn't separate policies",
         "Running the baselines instead of assuming they'd differ"],
        ["Pretrained model had zero dropout",
         "Asserting MC-dropout actually enables layers"],
        ["int8 model unexecutable in production",
         "Creating a session, not just checking numerical agreement"],
        ["ReLU zeroed 11 of 14 activation maps",
         "Noticing CAM is an evidence field, not a gradient product"],
        ["Rate limiter bucketed the whole world",
         "Asserting two clients behind one proxy get separate budgets"],
    ], y=2.4, widths=[5.9, 6.0], row_h=0.66, size=13)
    panel(s, 0.7, 5.6, 11.9, 1.0)
    text(s, 1.0, 5.8, 11.3, 0.7,
         "Every one was found by testing a claim the system makes about itself. "
         "None would have been caught by checking whether the app looked right.",
         size=14, colour=HIGHLIGHT)
    footer(s)

    # 15 — close
    s = blank(prs); eyebrow(s, "conclusion", 15, T)
    title(s, "Trust is earned by saying what you don't know.", y=2.3, size=40)
    text(s, 0.7, 3.6, 11.9, 1.2,
         "A clinical system earns trust not by being confident, but by being able "
         "to state — precisely, and verifiably — the limits of what it knows.",
         size=17, colour=MID, spacing=1.2)
    rule(s, 0.7, 5.0, 3.0, INSTRUMENT, 2.0)
    text(s, 0.7, 5.3, 11.9, 0.4, "github.com/krish2105/Deep-Learning-",
         size=14, colour=MID, font=MONO)
    text(s, 0.7, 5.85, 11.9, 0.4, TEAM, size=12, colour=MID)
    footer(s)

    path = OUT / "SENTINEL-CXR_Final.pptx"
    prs.save(path)
    return path, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    for fn in (build_week6, build_final):
        p, n = fn()
        print(f"Wrote {p.name}  ({n} slides)")
